import os
import io
import time
from dotenv import load_dotenv
from pypdf import PdfReader
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from google import genai
from pinecone import Pinecone, ServerlessSpec

load_dotenv()

pc = Pinecone(api_key=os.getenv("PINECONE_API_KEY"))
client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))

def get_embedding(text: str):
    result = client.models.embed_content(
        model="models/gemini-embedding-001",
        contents=text
    )
    return result.embeddings[0].values

def create_index_if_not_exists():
    index_name = os.getenv("PINECONE_INDEX_NAME")
    existing = [i.name for i in pc.list_indexes()]
    if index_name not in existing:
        pc.create_index(
            name=index_name,
            dimension=3072,
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1")
        )
    return pc.Index(index_name)

def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n[Slide {slide_num}]\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"
    return text

def chunk_text(text: str):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )
    return splitter.split_text(text)

def ingest_file(file_bytes: bytes, filename: str):
    index = create_index_if_not_exists()
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        text = extract_text_from_pdf(file_bytes)
    elif ext in ["pptx", "ppt"]:
        text = extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    chunks = chunk_text(text)
    vectors = []
    for i, chunk in enumerate(chunks):
        if i > 0 and i % 80 == 0:
            time.sleep(60)
        embedding = get_embedding(chunk)
        vectors.append({
            "id": f"{filename}-chunk-{i}",
            "values": embedding,
            "metadata": {"text": chunk, "source": filename}
        })
    index.upsert(vectors=vectors)
    return len(chunks)